"""
Handles PowerPoint (.pptx) file creation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
from slide_ai.layout import apply_background_color, get_layout_options
from slide_ai.image_editor import add_rounded_corners, pil_image_to_stream
import random

def create_pptx_with_unsplash(slide_data_list, topic, app_name="SlideAI", filename=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    CONTENT_LAYOUT = prs.slide_layouts[5]
    layout_opts = get_layout_options()

    # Title slide
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    apply_background_color(slide, (255, 255, 255))  # White bg for title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = topic
        tf = title_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(48)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Content slides
    for slide_data in slide_data_list:
        slide = prs.slides.add_slide(CONTENT_LAYOUT)
        # Pick a random or cycling background color
        bg_color = random.choice(layout_opts['background_colors'])
        apply_background_color(slide, bg_color)
        # Pick alignment randomly or by preference
        alignment = random.choice(layout_opts['alignments'])
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "Untitled Slide")
            tf = title_shape.text_frame
            if alignment == "left":
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif alignment == "center":
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif alignment == "right":
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(32)
            tf.vertical_anchor = MSO_ANCHOR.TOP
        # Body
        content = slide_data.get("content_points", [])
        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            for point in content:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(18)
            # Attribution
            if slide_data.get("unsplash_photographer_name"):
                p = tf.add_paragraph()
                p.text = f"Photo by {slide_data['unsplash_photographer_name']} on Unsplash"
                p.font.size = Pt(8)
                p.font.italic = True
        # Image
        image_stream = slide_data.get("actual_image_stream")
        if image_stream:
            image_stream.seek(0)
            pil_img = Image.open(image_stream)
            rounded_img = add_rounded_corners(pil_img, radius=60)
            rounded_stream = pil_image_to_stream(rounded_img)
            # Place image on right half, vertically centered
            left = Inches(7)
            top = Inches(2)
            width = Inches(5.5)
            slide.shapes.add_picture(rounded_stream, left, top, width=width)

    # Thank you slide
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    apply_background_color(slide, (255, 255, 255))
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Thank You"
        tf = title_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(54)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if filename is None:
        # Generate a default filename if none is provided
        filename = f"{topic.replace(' ', '_').lower()}_presentation.pptx"
    
    # Ensure the filename has the .pptx extension
    if not filename.lower().endswith('.pptx'):
        filename += '.pptx'
        
    prs.save(filename)
    return filename
